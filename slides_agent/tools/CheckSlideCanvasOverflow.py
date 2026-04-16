"""Check for content overflowing the slide canvas boundaries."""

import os
import subprocess
import tempfile
from copy import deepcopy
from pathlib import Path
from typing import List, Optional

from agency_swarm.tools import BaseTool
from pydantic import Field



class CheckSlideCanvasOverflow(BaseTool):
    """
    Detect slides with content overflowing the original canvas boundaries.

    This tool works by:
    1. Adding grey padding around each slide
    2. Rendering the padded presentation to images
    3. Inspecting the padding margins for non-grey pixels

    If content extends beyond the original slide boundaries, it will appear
    in the padding area and be detected as overflow.

    Returns a list of failing slide indices (1-based) and paths to debug
    images showing the overflow.

    Requires: LibreOffice (soffice) and pdf2image Python package.
    """

    input_pptx: str = Field(
        ...,
        description="Path to the input PowerPoint file (.pptx)",
    )
    max_width_px: int = Field(
        default=1600,
        description="Maximum width in pixels for rendered images (default 1600)",
    )
    max_height_px: int = Field(
        default=900,
        description="Maximum height in pixels for rendered images (default 900)",
    )
    pad_px: int = Field(
        default=100,
        description="Padding in pixels to add on each side (default 100)",
    )

    def run(self) -> str:
        """Check for canvas overflow and return results."""
        input_path = Path(self.input_pptx)

        # Validate input
        if not input_path.exists():
            return f"Error: Input file not found: {self.input_pptx}"
        if input_path.suffix.lower() != ".pptx":
            return f"Error: Input must be a PowerPoint file (.pptx), got: {input_path.suffix}"

        # Check for required external tools
        if not self._check_soffice():
            return "Error: LibreOffice (soffice) not found. Please install LibreOffice."

        try:
            import numpy as np
            from pdf2image import convert_from_path
            from PIL import Image
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.util import Emu
        except ImportError as e:
            return f"Error: Missing required package: {e}"

        # Import render_slides for DPI calculation
        try:
            from .render_slides import calc_dpi_via_ooxml, rasterize
        except ImportError:
            return "Error: Could not import render_slides module"

        # Constants
        EMU_PER_INCH = 914_400
        PAD_RGB = (200, 200, 200)

        def px_to_emu(px: int, dpi: int) -> Emu:
            return Emu(int(px * EMU_PER_INCH // dpi))

        # Calculate DPI
        try:
            dpi = calc_dpi_via_ooxml(
                str(input_path), self.max_width_px, self.max_height_px
            )
        except Exception as e:
            return f"Error calculating DPI: {e}"

        # Create temporary directory for work
        tmpdir = tempfile.mkdtemp(prefix="overflow_check_")
        enlarged_pptx = Path(tmpdir) / "enlarged.pptx"

        try:
            # Enlarge the deck with padding
            prs = Presentation(str(input_path))
            w0 = prs.slide_width
            h0 = prs.slide_height
            pad_emu = px_to_emu(self.pad_px, dpi)
            w1 = Emu(w0 + 2 * pad_emu)
            h1 = Emu(h0 + 2 * pad_emu)
            prs.slide_width = w1
            prs.slide_height = h1

            for slide in prs.slides:
                # Shift all shapes
                for shp in list(slide.shapes):
                    shp.left = Emu(int(shp.left) + pad_emu)
                    shp.top = Emu(int(shp.top) + pad_emu)

                # Add padding rectangles
                pads = (
                    (Emu(0), Emu(0), pad_emu, h1),  # left
                    (Emu(int(w1) - int(pad_emu)), Emu(0), pad_emu, h1),  # right
                    (Emu(0), Emu(0), w1, pad_emu),  # top
                    (Emu(0), Emu(int(h1) - int(pad_emu)), w1, pad_emu),  # bottom
                )

                sp_tree = slide.shapes._spTree

                for left, top, width, height in pads:
                    pad_shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
                    )
                    pad_shape.fill.solid()
                    pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
                    pad_shape.line.fill.background()
                    sp_tree.remove(pad_shape._element)
                    sp_tree.insert(2, pad_shape._element)

            prs.save(str(enlarged_pptx))

            # Render to images
            img_paths = rasterize(str(enlarged_pptx), str(Path(tmpdir) / "imgs"), dpi)

            # Calculate padding ratios
            pad_ratio_w = pad_emu / w1
            pad_ratio_h = pad_emu / h1

            # Inspect images for overflow
            tol = max(1, round((300 - dpi) / 25)) if dpi < 300 else 0
            tol = min(tol, 10)
            pad_colour = np.array(PAD_RGB, dtype=np.uint8)
            failures = []

            for idx, img_path in enumerate(img_paths, start=1):
                with Image.open(img_path) as img:
                    rgb = img.convert("RGB")
                    arr = np.asarray(rgb)

                h, w, _ = arr.shape
                pad_x = int(w * pad_ratio_w) - 1
                pad_y = int(h * pad_ratio_h) - 1

                margins = [
                    arr[:, :pad_x, :],  # left
                    arr[:, w - pad_x:, :],  # right
                    arr[:pad_y, :, :],  # top
                    arr[h - pad_y:, :, :],  # bottom
                ]

                def is_clean(margin):
                    diff = np.abs(margin.astype(np.int16) - pad_colour)
                    matches = np.all(diff <= tol, axis=-1)
                    mismatch_fraction = 1.0 - (np.count_nonzero(matches) / matches.size)
                    max_mismatch = 0.01 if dpi >= 300 else (0.02 if dpi >= 200 else 0.03)
                    return mismatch_fraction <= max_mismatch

                if not all(is_clean(m) for m in margins):
                    failures.append((idx, img_path))

            if failures:
                result = f"OVERFLOW DETECTED on {len(failures)} slide(s):\n"
                for idx, img_path in failures:
                    result += f"  - Slide {idx}: {img_path}\n"
                result += "\nDebug images show grey padding with overflow content visible."
                return result
            else:
                return "No overflow detected. All slides are within canvas boundaries."

        except Exception as e:
            return f"Error checking overflow: {e}"

    def _check_soffice(self) -> bool:
        """Check if LibreOffice is available."""
        try:
            kwargs = {
                "capture_output": True,
                "timeout": 15,
                "text": True,
                "input": "\n",
            }
            if os.name == "nt":
                kwargs["creationflags"] = 0x08000000  # CREATE_NO_WINDOW

            soffice_bin = "soffice.com" if os.name == "nt" else "soffice"
            subprocess.run([soffice_bin, "--version"], **kwargs)
            return True
        except (FileNotFoundError, subprocess.TimeoutExpired):
            return False


if __name__ == "__main__":
    # Test with a sample file if available
    test_pptx = Path(__file__).parent.parent / "files" / "test.pptx"
    if test_pptx.exists():
        tool = CheckSlideCanvasOverflow(input_pptx=str(test_pptx))
        print(tool.run())
    else:
        print(f"Test file not found: {test_pptx}")
        print("Tool definition is valid.")
